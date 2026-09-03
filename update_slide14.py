import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TARGET_FILES = [
    "/Users/Harika/Downloads/Smart-Horizon-2026-48Hour-Intrnl Hackathon Grand Finale Template (1).pptx",
    "/Users/Harika/Downloads/Smart-Horizon-2026-48Hour-Intrnl Hackathon Grand Finale Template.pptx",
    "/Users/Harika/Desktop/Smart-Horizon-2026-48Hour-Intrnl Hackathon Grand Finale Template (1).pptx",
    "/Users/Harika/Desktop/SUTRA/SUTRA/Smart-Horizon-2026-48Hour-Intrnl Hackathon Grand Finale Template (1).pptx",
    "/Users/Harika/Desktop/SUTRA/SUTRA/Smart_Horizon_2026_SUTRA_Grand_Finale_Pitch.pptx",
    "/Users/Harika/Downloads/Smart_Horizon_2026_SUTRA_Grand_Finale_Pitch.pptx"
]

COLOR_PRIMARY = RGBColor(15, 23, 42)      # Deep Navy / Slate
COLOR_HEADER = RGBColor(14, 116, 144)     # Cyan / Teal Blue
COLOR_TEXT = RGBColor(30, 41, 59)         # Dark Charcoal Body
COLOR_MUTED = RGBColor(71, 85, 105)       # Slate Muted

for file_path in TARGET_FILES:
    if not os.path.exists(file_path):
        continue
        
    prs = Presentation(file_path)
    s14 = prs.slides[13]
    
    title_shape = None
    info_box = None
    
    for sh in list(s14.shapes):
        if sh.name == "Title 1" or (sh.has_text_frame and "Thank You" in sh.text_frame.text):
            title_shape = sh
        elif sh.has_text_frame and sh.name not in ["Date Placeholder 2", "Slide Number Placeholder 4", "Footer Placeholder 7"]:
            info_box = sh

    # 1. Position "Thank You" right at the TOP
    if title_shape:
        title_shape.left = Inches(0.50)
        title_shape.top = Inches(0.40)
        title_shape.width = Inches(9.00)
        title_shape.height = Inches(0.90)
        title_shape.text_frame.clear()
        p = title_shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = "Thank You"
        r.font.name = "Times New Roman"
        r.font.size = Pt(40)
        r.font.bold = True
        r.font.color.rgb = COLOR_PRIMARY
        
    # 2. Position all other information LOWER DOWN (leaving clean breathing space)
    if not info_box:
        info_box = s14.shapes.add_textbox(Inches(0.50), Inches(1.95), Inches(9.00), Inches(3.05))
    else:
        info_box.left = Inches(0.50)
        info_box.top = Inches(1.95)
        info_box.width = Inches(9.00)
        info_box.height = Inches(3.05)
        
    tf = info_box.text_frame
    tf.word_wrap = True
    tf.clear()
    
    # Project Title
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    p0.space_after = Pt(4)
    r0 = p0.add_run()
    r0.text = "Project SUTRA — Swarm Unified Tactical Reconnaissance Architecture"
    r0.font.name = "Times New Roman"
    r0.font.size = Pt(16)
    r0.font.bold = True
    r0.font.color.rgb = COLOR_HEADER
    
    # Track & College Details
    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(6)
    p1.space_after = Pt(12)
    p1.line_spacing = 1.2
    r1 = p1.add_run()
    r1.text = "Problem Statement: SH-DST-05  |  Track: Defence & SpaceTech (DST)\nTeam ID: SHIH26-TID-361  |  New Horizon College of Engineering (NHCE), Bengaluru"
    r1.font.name = "Times New Roman"
    r1.font.size = Pt(12)
    r1.font.color.rgb = COLOR_TEXT
    
    # Team Members Header
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)
    p2.space_after = Pt(3)
    r2 = p2.add_run()
    r2.text = "Grand Finals Project Team:"
    r2.font.name = "Times New Roman"
    r2.font.size = Pt(12)
    r2.font.bold = True
    r2.font.color.rgb = COLOR_PRIMARY
    
    # Team Members List
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_after = Pt(16)
    r3 = p3.add_run()
    r3.text = "Nikhil (Team Lead)  •  Vedanth Sai Ram  •  Siva Kesava  •  Harika  •  Rohith Kumar"
    r3.font.name = "Times New Roman"
    r3.font.size = Pt(12.5)
    r3.font.bold = True
    r3.font.color.rgb = COLOR_HEADER
    
    # Invitation to Demonstration
    p4 = tf.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(8)
    p4.space_after = Pt(4)
    r4 = p4.add_run()
    r4.text = "We welcome questions from the jury and warmly invite you to inspect our\nlive Gazebo Sim 8 digital twin, WebGPU GCS & 232-test verification suite."
    r4.font.name = "Times New Roman"
    r4.font.size = Pt(11.5)
    r4.font.italic = True
    r4.font.color.rgb = COLOR_MUTED
    
    prs.save(file_path)
    print(f"Saved: {file_path}")

