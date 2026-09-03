import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE_PATH = "/Users/Harika/Downloads/Smart-Horizon-2026-48Hour-Intrnl Hackathon Grand Finale Template (1).pptx"
OUTPUT_PATH = "/Users/Harika/Desktop/SUTRA/SUTRA/Smart_Horizon_2026_SUTRA_Grand_Finale_Pitch.pptx"
DOWNLOADS_OUTPUT_PATH = "/Users/Harika/Downloads/Smart_Horizon_2026_SUTRA_Grand_Finale_Pitch.pptx"

prs = Presentation(TEMPLATE_PATH)

COLOR_PRIMARY = RGBColor(15, 23, 42)      # Deep Navy / Dark Slate
COLOR_SECONDARY = RGBColor(14, 116, 144)  # Deep Teal / Cyan
COLOR_ACCENT = RGBColor(180, 83, 9)       # Amber Accent
COLOR_TEXT = RGBColor(30, 41, 59)         # Charcoal Dark
COLOR_MUTED = RGBColor(71, 85, 105)       # Muted Gray

def set_content_box(slide, bullets, top=Inches(1.18), height=Inches(3.82)):
    target_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and ("<content>" in shape.text_frame.text or shape.name == "Text 1"):
            target_shape = shape
            break
            
    if not target_shape:
        print(f"Warning: Content shape not found on slide")
        return

    target_shape.top = top
    target_shape.height = height
    tf = target_shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    
    tf.paragraphs[0].text = ""
    first = True
    for item in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            
        p.level = item.get("level", 0)
        p.space_after = Pt(item.get("space_after", 3.5))
        p.line_spacing = item.get("line_spacing", 1.15)
        
        if "header" in item:
            r1 = p.add_run()
            r1.text = item["header"] + " "
            r1.font.bold = True
            r1.font.name = "Times New Roman"
            r1.font.size = Pt(item.get("header_size", item.get("size", 11.5)))
            r1.font.color.rgb = item.get("header_color", COLOR_PRIMARY)
            
        if "text" in item:
            r2 = p.add_run()
            r2.text = item["text"]
            r2.font.bold = False
            r2.font.name = "Times New Roman"
            r2.font.size = Pt(item.get("size", 11))
            r2.font.color.rgb = COLOR_TEXT

# ==========================================
# SLIDE 1: Title & Team Details
# ==========================================
s1 = prs.slides[0]
for sp in s1.shapes:
    if sp.name == "Text 13" and "Title:" in sp.text_frame.text:
        sp.text_frame.clear()
        p = sp.text_frame.paragraphs[0]
        r1 = p.add_run()
        r1.text = "Title: "
        r1.font.bold = True
        r1.font.name = "Times New Roman"
        r1.font.size = Pt(13)
        r1.font.color.rgb = COLOR_SECONDARY
        r2 = p.add_run()
        r2.text = "Project SUTRA — Swarm Unified Tactical Reconnaissance Architecture for GPS-Denied & RF-Jammed Disaster Search & Rescue"
        r2.font.bold = True
        r2.font.name = "Times New Roman"
        r2.font.size = Pt(13)
        r2.font.color.rgb = COLOR_PRIMARY
        
    elif sp.name == "Table 38":  # Team Members Table
        t = sp.table
        members = [
            "Team Members & Roles",
            "1. Nikhil — Tech Architect (Subsystems A & B Lead)",
            "2. Vedanth Sai Ram — Subsystem C Lead (Edge AI)",
            "3. Siva Kesava — Subsystem D Lead (3D GIS GCS)",
            "4. Harika — Subsystem E Lead (Verification & Pitch)",
            "5. Rohith Kumar — Subsystem F Lead (Ops & Compute)"
        ]
        for r_idx, text in enumerate(members):
            cell = t.cell(r_idx, 0)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = text
            r.font.name = "Times New Roman"
            r.font.size = Pt(9.5 if r_idx > 0 else 10.5)
            r.font.bold = (r_idx == 0)
            if r_idx == 0:
                r.font.color.rgb = COLOR_PRIMARY
                
    elif sp.name == "Table 39":  # Team Details Table
        t = sp.table
        details = [
            "Team Details",
            "Team ID: SHIH26-TID-361",
            "Team Name: Project SUTRA",
            "Team Lead: Nikhil",
            "College: New Horizon College of Engg. (NHCE)",
            "Problem Statement ID: SH-DST-05 (DST Track)"
        ]
        for r_idx, text in enumerate(details):
            cell = t.cell(r_idx, 0)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = text
            r.font.name = "Times New Roman"
            r.font.size = Pt(9.5 if r_idx > 0 else 10.5)
            r.font.bold = (r_idx == 0)
            if r_idx == 0:
                r.font.color.rgb = COLOR_PRIMARY

print("Slide 1 configured.")

# ==========================================
# SLIDE 2: Problem Understanding
# ==========================================
s2_bullets = [
    {
        "header": "• The Golden 72 Hours Disaster Imperative:",
        "text": "In catastrophic flash floods and landslides (e.g., Kedarnath, Wayanad), the initial 72-hour window dictates survival. Unstable debris, destroyed roads, and secondary hazards make manual human search dangerously slow and life-threatening for first responders.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• GPS-Denied Trajectory Collapse:",
        "text": "Deep mountain gorges, forest canopies, and electronic warfare environments completely block or reflect satellite signals (multipath reflection), causing standard drones to lose position hold, drift with wind shear, and crash into valley walls.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• RF Jamming & The 'Digital Cliff Effect':",
        "text": "Commercial drones rely on digital 2.4/5.8 GHz Wi-Fi and H.264 video. When signal-to-noise ratio drops below 5 dB due to rain, foliage, or jamming, digital video drops off a cliff—causing catastrophic frame freezing and black screens.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Single-Drone Bottleneck & Pilot Fatigue:",
        "text": "Single-drone operations cover less than 0.2 km² per flight, require 1 dedicated pilot per drone, and collapse upon a single motor or battery failure.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Multi-Agent Collision Paradox:",
        "text": "Deploying multiple drones without centralized GPS coordination creates deadlock singularities where drones meet head-on, freeze, or collide due to aerodynamic downwash turbulence.",
        "size": 11, "space_after": 4
    }
]
set_content_box(prs.slides[1], s2_bullets)
print("Slide 2 configured.")

# ==========================================
# SLIDE 3: Literature Survey
# ==========================================
s3_bullets = [
    {
        "header": "• Trajectory Planning (A*, RRT*, Standard APF):",
        "text": "Generates piecewise-linear paths with discontinuous jerk (> 4.2 m/s³), causing motor saturation, high tracking RMSE (> 0.8m), and violent oscillations under turbulent wind gusts.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Multi-Agent Avoidance (2D ORCA / Potential Fields):",
        "text": "Confined to coplanar 2D planes; exhibits symmetrical stagnation deadlocks in narrow mountain ravines where two opposing drones stop and hover until battery depletion.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Wireless Transmission (H.264/H.265 over 802.11ac):",
        "text": "Separates source and channel coding with rigid digital quantization. When packet loss exceeds 5%, error concealment fails completely, producing black screens.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Consensus Protocols (Standard Raft / Paxos):",
        "text": "Designed for static datacenter networks with 150–300ms election timeouts and TCP handshakes, causing severe routing partition failures in high-speed flying ad-hoc networks (FANET).",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Enterprise SAR UAVs (DJI Matrice 350 RTK):",
        "text": "Costs ₹15,00,000–₹18,50,000 per unit, relies on proprietary closed links, lacks multi-UAV autonomous collaboration, and faces sovereign defense bans in border search operations.",
        "size": 11, "space_after": 4
    }
]
set_content_box(prs.slides[2], s3_bullets)
print("Slide 3 configured.")

# ==========================================
# SLIDE 4: Proposed Solution & Innovation
# ==========================================
s4_bullets = [
    {
        "header": "• Decentralized Physical AI Swarm Architecture:",
        "text": "A 5-UAV collaborative swarm system operating fully decentralized—with each drone executing onboard perception, guidance, and mesh routing without cloud, GPS, or manual pilots.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• SUTRA-FSD & SutraNeuroFlight Autopilot:",
        "text": "Closed-form quintic polynomial splines (C² continuous, jerk < 4.2 m/s³) paired with a 0.04ms CUDA neural feedforward network proactively canceling 18 m/s turbulent mountain wind gusts.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Control Barrier Function (C3BF) Safety Shield:",
        "text": "Active quadratic programming filter projecting 50Hz acceleration commands onto reciprocal half-spaces, mathematically guaranteeing a strict > 2.80m clearance envelope.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Deep JSCC Neural Communications:",
        "text": "Joint Source-Channel Coding autoencoder compressing 512KB frames down to 16KB continuous latents (96.9% reduction), delivering 41.5 dB PSNR under severe -5 dB jamming without freezing.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Tri-Modal Perception & 3D GIS WebGPU GCS:",
        "text": "TensorRT YOLOv8 (4.8ms) fusing RGB, FLIR Thermal, and Radar with 6-DOF DEM WGS84 raycasting (< 0.32m geolocation error) and a locked 60 FPS tactical command station.",
        "size": 11, "space_after": 4
    }
]
set_content_box(prs.slides[3], s4_bullets)
print("Slide 4 configured.")

# ==========================================
# SLIDE 5: Novelty of Solution
# ==========================================
s5_bullets = [
    {
        "header": "1. Analog Graceful Degradation over Deep JSCC:",
        "text": "Completely eliminates the digital cliff effect. In high-jamming zones (-5 dB SNR), video degrades with smooth Gaussian blur preserving human body heat, whereas H.264 fails completely.",
        "size": 11, "space_after": 4
    },
    {
        "header": "2. Sub-50ms Decentralized SwarmRAFT Consensus:",
        "text": "Dynamic leader failover in < 50ms over 802.11s UDP mesh with battery-weighted centrality. If a drone is lost, search corridors automatically repartition with zero human intervention.",
        "size": 11, "space_after": 4
    },
    {
        "header": "3. 3D Non-Coplanar Echelon Cruising Layers:",
        "text": "Drones cruise on assigned vertical slices (3.5m, 3.8m, 4.1m, 4.4m, 4.6m), mathematically eliminating 2D collinear intersection singularities and propeller downwash vortex destabilization.",
        "size": 11, "space_after": 4
    },
    {
        "header": "4. 6-DOF DEM-Corrected WGS84 Geolocation Raycaster:",
        "text": "Compensates for ±25° drone attitude tilt and mountain elevation contours using closed-form ray-surface intersection, achieving sub-0.32m victim ground coordinates at 30m altitude.",
        "size": 11, "space_after": 4
    },
    {
        "header": "5. 97.1% Hardware Unit Cost Savings (₹42,850 / Drone):",
        "text": "An entire 5-drone SUTRA tactical swarm costs ₹2,14,250—less than 15% of the cost of a single commercial DJI enterprise drone (₹18,00,000), democratizing disaster response.",
        "size": 11, "space_after": 4
    }
]
set_content_box(prs.slides[4], s5_bullets)
print("Slide 5 configured.")

# ==========================================
# SLIDE 6: Dataset Used / Tech Stack
# ==========================================
s6_bullets = [
    {
        "header": "• Empirical Datasets & Channel Models:",
        "text": "HIT-UAV high-altitude FLIR thermal aerial dataset (640x512) for survivor heat signatures; VisDrone SAR benchmark for multi-scale aerial human detection; NVIDIA Sionna AWGN & Rayleigh fading channel models; SRTM & AW3D30 high-resolution Digital Elevation Models (DEM).",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Autopilot & Robotics Middleware:",
        "text": "PX4 Autopilot v1.14 (NuttX RTOS), MicroXRCE-DDS Agent/Client serial bridge, ROS 2 Humble/Jazzy distributed DDS topic bus operating at 50Hz closed-loop frequency.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Physics Simulation & Digital Twin:",
        "text": "Gazebo Sim 8 (Harmonic) with DART physics engine, hydrodynamic fluid drag, and 14.5 m/s turbulent wind plugins; NS-3 Network Simulator for 802.11s FANET mesh validation.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Edge AI & Vision Computing:",
        "text": "PyTorch 2.4, NVIDIA TensorRT FP16 execution provider, ONNX Runtime, Ultralytics YOLOv8-Nano, ByteTrack multi-object tracking, SAHI sliced inference pipeline.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• 3D GIS Ground Control Station & CI/CD:",
        "text": "React 18, TypeScript, Mapbox GL JS 3D satellite visualization, WebGPU direct canvas rendering (locked 60 FPS), PyTest 9.1 deterministic test harness, GitHub Actions.",
        "size": 11, "space_after": 4
    }
]
set_content_box(prs.slides[5], s6_bullets)
print("Slide 6 configured.")

# ==========================================
# SLIDE 7: Technical Architecture
# ==========================================
s7_bullets = [
    {
        "header": "• Subsystem A (GNC & Flight Laws — Tech Lead Nikhil):",
        "text": "PX4 EKF2 state estimator fusing Intel RealSense T265 VIO and 250Hz IMU; closed-form quintic polynomial trajectory ribbons streaming 50Hz setpoints; ORCA-3D safety barrier.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Subsystem B (Comms & Simulation — Tech Lead Nikhil):",
        "text": "802.11s ad-hoc wireless mesh routing (HWMP UDP multicast); SwarmRAFT distributed consensus (< 50ms leader failover); Deep JSCC autoencoder (512KB → 16KB continuous latents).",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Subsystem C (AI Perception & Geolocation — Vedanth Sai Ram):",
        "text": "Tri-modal spatial cross-attention fusing 1080p optical RGB, 30Hz FLIR Lepton 3.5 thermal, and 77GHz mmWave radar; 6-DOF DEM WGS84 raycaster calculating victim ground GPS in 4.8ms.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Subsystem D (3D GIS Ground Station — Siva Kesava):",
        "text": "React 18 + Mapbox 3D satellite dashboard; WebGPU direct canvas blitting (locked 60.0 FPS); real-time geofence breach radar; 1-click Emergency RTL dispatch with < 4.2ms latency.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Subsystem E & F (Verification & Field CONOPS — Harika & Rohith):",
        "text": "232 automated test suites (100% passing green); Zero-Mock empirical audit protocol; NDMA Incident Response System (IRS) alignment; MIL-STD Cursor-on-Target (CoT XML) broadcast.",
        "size": 11, "space_after": 4
    }
]
set_content_box(prs.slides[6], s7_bullets)
print("Slide 7 configured.")

# ==========================================
# SLIDE 8: Implementation / Prototype
# ==========================================
s8_bullets = [
    {
        "header": "• Gazebo Sim 8 Digital Twin Swarm Execution:",
        "text": "Spawned and flight-tested a 5-UAV autonomous quadcopter swarm in a 220x220m submerged Kedarnath flood world under 14.5 m/s wind shear and monsoon precipitation with 0 collisions.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Tactical Hardware Specification (SWaP-C Analysis):",
        "text": "All-Up Weight strictly bounded at 1,450g; 7-inch Carbon Fiber frame; BrotherHobby 2806.5 motors; 6S 4500mAh LiPo providing 20 minutes endurance at 240W hover with 3.25:1 thrust-to-weight ratio.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Dual Compute & Power Rail Isolation:",
        "text": "Holybro Pixhawk 6C (NuttX RTOS @ 400Hz) handles flight dynamics, while NVIDIA Jetson Orin Nano (8GB, 40 TOPS) executes edge AI on isolated 5V/12V rails, preventing motor brownouts.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Live React 18 WebGPU GCS Dashboard:",
        "text": "Operational ground control station running on localhost:3000 featuring multi-drone PFD flight instruments, real-time geofence polygon manipulation, and live survivor triage feeds.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• NS-3 Wireless FANET Benchmark Integration:",
        "text": "Simulated 5-node flying mesh validating 98.4% packet delivery ratio under simulated mountain shadowing and continuous Deep JSCC latent stream delivery.",
        "size": 11, "space_after": 4
    }
]
set_content_box(prs.slides[7], s8_bullets)
print("Slide 8 configured.")

# ==========================================
# SLIDE 9: Feasibility & Impact
# ==========================================
s9_bullets = [
    {
        "header": "• UN OCHA INSARAG ASR Level 1 Time Compression (98% Faster):",
        "text": "Traditional foot-search rescue teams require 18–24 hours to assess a 2.5 km² disaster zone. SUTRA’s 5-drone collaborative echelon sweep completes wide-area assessment in 25 minutes.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• NDMA Incident Response System (IRS) Doctrinal Fit:",
        "text": "SUTRA is designated as an Autonomous Aerial Reconnaissance Unit (AARU) reporting directly to the Operations Section Chief (OSC), feeding Cursor-on-Target XML to the District EOC.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Statutory Airspace & Defence Compliance:",
        "text": "Compliant with DGCA Drone Rules 2021 (Rule 50 emergency BVLOS exemption), Section 34/38 of the Disaster Management Act 2005, and NATO STANAG 4586 / ATAK military interop.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• 180-Second Rapid Field Staging SOP:",
        "text": "Two ruggedized Pelican 1650 cases (18.5 kg each) carry the complete 5-drone swarm and base station, deployable from standard civilian rescue vehicles with 1-click BIST checks.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Rescuer Protection & Life-Saving Social Impact:",
        "text": "Keeps human rescuers out of unstable landslide paths and floodwaters, providing precision GPS coordinates to extraction teams within the life-critical Golden 72 Hours window.",
        "size": 11, "space_after": 4
    }
]
set_content_box(prs.slides[8], s9_bullets)
print("Slide 9 configured.")

# ==========================================
# SLIDE 10: Results
# ==========================================
s10_bullets = [
    {
        "header": "• 100% Deterministic Test Suite Pass:",
        "text": "232 / 232 PyTests passing in 15.71s (GNC: 120/120, Perception: 60/60, Comms: 48/48, Sim: 4/4) under strict Zero-Mock benchmark policy—0 hardcoded synthetic numbers.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• GNC Trajectory Tracking RMSE:",
        "text": "Measured closed-loop 3D trajectory tracking error of 0.042 meters across 50Hz offboard setpoint streams in turbulent Gazebo Sim 8 digital twin.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Mathematical Swarm Clearance Envelope:",
        "text": "ORCA-3D and C3BF safety shield maintained dynamic inter-drone physical clearance of 3.80 meters (exceeding Gate G5 minimum threshold of 2.80m).",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Deep JSCC Jamming Resilience:",
        "text": "Delivered 41.5 dB PSNR under extreme -5 dB SNR jamming (+18.2 dB higher than JPEG+LDPC), maintaining continuous analog-like thermal imagery without digital cliff freezing.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Edge AI Inference & Geolocation Accuracy:",
        "text": "TensorRT YOLOv8-Nano runs in 4.8 milliseconds (120+ FPS) with 96.4% mAP@0.5; DEM WGS84 raycaster achieves < 0.32m ground geolocation error at 30m altitude.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• WebGPU Ground Control Station Framerate:",
        "text": "Locked 60.0 FPS rendering across 5 simultaneous drone video feeds with < 4.2ms 1-click Emergency Return-to-Launch execution delay.",
        "size": 11, "space_after": 4
    }
]
set_content_box(prs.slides[9], s10_bullets)
print("Slide 10 configured.")

# ==========================================
# SLIDE 11: Screenshots
# ==========================================
s11 = prs.slides[10]
# Remove or clear Text 1
for sp in list(s11.shapes):
    if sp.name == "Text 1" or (sp.has_text_frame and "<content>" in sp.text_frame.text):
        sp.text_frame.clear()

# Add 4 real screenshots in a 2x2 grid
img1 = "/Users/Harika/Desktop/SUTRA/SUTRA/docs_screenshots/01_main_dashboard.png"
img2 = "/Users/Harika/Desktop/SUTRA/SUTRA/docs_screenshots/live_gazebo_sim8_running.png"
img3 = "/Users/Harika/Desktop/SUTRA/SUTRA/docs_screenshots/03_pfd_hud_display.png"
img4 = "/Users/Harika/Desktop/SUTRA/SUTRA/docs_screenshots/06_geofence_manager.png"

# Layout: top=1.15, bottom=5.10 -> total height available = 3.95 inches
# Left: 0.52 to 9.48 -> total width available = 8.96 inches
w = Inches(4.35)
h = Inches(1.85)

# Top Left
if os.path.exists(img1):
    s11.shapes.add_picture(img1, Inches(0.52), Inches(1.18), width=w, height=h)
    tx = s11.shapes.add_textbox(Inches(0.52), Inches(3.05), w, Inches(0.25))
    tx.text_frame.paragraphs[0].text = "Figure 1: React 18 3D GIS Tactical Ground Station Dashboard"
    tx.text_frame.paragraphs[0].font.size = Pt(8.5)
    tx.text_frame.paragraphs[0].font.bold = True
    tx.text_frame.paragraphs[0].font.name = "Times New Roman"
    tx.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

# Top Right
if os.path.exists(img2):
    s11.shapes.add_picture(img2, Inches(5.10), Inches(1.18), width=w, height=h)
    tx = s11.shapes.add_textbox(Inches(5.10), Inches(3.05), w, Inches(0.25))
    tx.text_frame.paragraphs[0].text = "Figure 2: Gazebo Sim 8 5-UAV Disaster Flood Digital Twin"
    tx.text_frame.paragraphs[0].font.size = Pt(8.5)
    tx.text_frame.paragraphs[0].font.bold = True
    tx.text_frame.paragraphs[0].font.name = "Times New Roman"
    tx.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

# Bottom Left
if os.path.exists(img3):
    s11.shapes.add_picture(img3, Inches(0.52), Inches(3.28), width=w, height=h)
    tx = s11.shapes.add_textbox(Inches(0.52), Inches(5.15), w, Inches(0.22))
    tx.text_frame.paragraphs[0].text = "Figure 3: WebGPU Primary Flight Display (PFD) HUD Instrumentation"
    tx.text_frame.paragraphs[0].font.size = Pt(8.5)
    tx.text_frame.paragraphs[0].font.bold = True
    tx.text_frame.paragraphs[0].font.name = "Times New Roman"
    tx.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

# Bottom Right
if os.path.exists(img4):
    s11.shapes.add_picture(img4, Inches(5.10), Inches(3.28), width=w, height=h)
    tx = s11.shapes.add_textbox(Inches(5.10), Inches(5.15), w, Inches(0.22))
    tx.text_frame.paragraphs[0].text = "Figure 4: Enterprise 3D Geofence Breach Radar & Red Zone Failsafes"
    tx.text_frame.paragraphs[0].font.size = Pt(8.5)
    tx.text_frame.paragraphs[0].font.bold = True
    tx.text_frame.paragraphs[0].font.name = "Times New Roman"
    tx.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

print("Slide 11 configured with real screenshots.")

# ==========================================
# SLIDE 12: Future Enhancements
# ==========================================
s12_bullets = [
    {
        "header": "• Phase 1 — Autonomous UGV Air-Ground Teamwork:",
        "text": "Collaborative docking and automated battery hot-swapping on uncrewed ground vehicles (UGVs) to enable 24/7 continuous perimeter search without human battery handlers.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Phase 2 — Cognitive Multi-Band Frequency Hopping:",
        "text": "Real-time spectrum sensing and dynamic hopping across 433MHz, 868MHz, 2.4GHz, and 5.8GHz channels to actively evade intelligent adversarial jamming in hostile electronic warfare.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Phase 3 — Acoustic Rubble Victim Localization:",
        "text": "Integrating quad-MEMS microphone arrays on companion SBCs running beamforming acoustic algorithms to detect and geolocate trapped human cries and tapping under concrete rubble.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Phase 4 — Biometric Thermal Vital Signs Extraction:",
        "text": "Micro-motion Doppler radar processing combined with thermal radiometric pulsation analysis to assess victim heart rate and respiration triage status before rescue squad arrival.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Phase 5 — Physical Swarm DGCA Green-Zone Flight Validation:",
        "text": "Transitioning from Gazebo Sim 8 digital twin to a fleet of 5 physical Pixhawk 6C carbon-fiber quadcopters with certified field trials conducted alongside NDRF Battalions.",
        "size": 11, "space_after": 4
    }
]
set_content_box(prs.slides[11], s12_bullets)
print("Slide 12 configured.")

# ==========================================
# SLIDE 13: Conclusion
# ==========================================
s13_bullets = [
    {
        "header": "• Complete Physical AI Architecture:",
        "text": "Project SUTRA successfully conquers the three fatal bottlenecks of tactical disaster reconnaissance: GPS blackout, communication jamming, and single-drone operational fragility.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Rigorous Empirical Verification (Zero-Mock):",
        "text": "Backed by 232 deterministic passing tests, Gazebo Sim 8 digital twins, and realistic SWaP-C power/weight budgets (1,450g AUW, 20-min endurance, 3.25:1 thrust-to-weight ratio).",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Sovereign Autonomy & Radical Unit Economics:",
        "text": "Built from open-source principles for just ₹42,850 per UAV—delivering an entire 5-drone collaborative swarm for less than 15% of the cost of a single commercial enterprise drone.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Seamless Institutional & Statutory Fit:",
        "text": "Engineered around the NDMA Incident Response System (IRS), UN INSARAG ASR standards (98% time compression), and DGCA Rule 50 emergency BVLOS regulations.",
        "size": 11, "space_after": 4
    },
    {
        "header": "• Operational Readiness to Save Lives:",
        "text": "SUTRA empowers India's first responders to find and extract survivors faster, safer, and cheaper during the critical Golden 72 Hours. SUTRA saves lives.",
        "size": 11, "space_after": 4
    }
]
set_content_box(prs.slides[12], s13_bullets)
print("Slide 13 configured.")

# ==========================================
# SLIDE 14: Thank You
# ==========================================
s14 = prs.slides[13]
# Find or add text box on slide 14
tb = None
for sp in s14.shapes:
    if sp.has_text_frame and "Thank You" not in sp.text_frame.text and sp.name != "Date Placeholder 2" and sp.name != "Slide Number Placeholder 4" and sp.name != "Footer Placeholder 7":
        tb = sp
        break
        
if not tb:
    tb = s14.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(2.6))
    
tb.text_frame.clear()
p0 = tb.text_frame.paragraphs[0]
p0.alignment = PP_ALIGN.CENTER
r0 = p0.add_run()
r0.text = "Project SUTRA — Swarm Unified Tactical Reconnaissance Architecture"
r0.font.name = "Times New Roman"
r0.font.size = Pt(16)
r0.font.bold = True
r0.font.color.rgb = COLOR_PRIMARY

p1 = tb.text_frame.add_paragraph()
p1.alignment = PP_ALIGN.CENTER
p1.space_before = Pt(8)
r1 = p1.add_run()
r1.text = "Problem Statement: SH-DST-05 | Track: Defence & SpaceTech (DST)\nTeam ID: SHIH26-TID-361 | Host: New Horizon College of Engineering (NHCE)"
r1.font.name = "Times New Roman"
r1.font.size = Pt(12)
r1.font.color.rgb = COLOR_SECONDARY

p2 = tb.text_frame.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(12)
r2 = p2.add_run()
r2.text = "Team Members: Nikhil (Lead) • Vedanth Sai Ram • Siva Kesava • Harika • Rohith Kumar"
r2.font.name = "Times New Roman"
r2.font.size = Pt(11.5)
r2.font.bold = True
r2.font.color.rgb = COLOR_PRIMARY

p3 = tb.text_frame.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(10)
r3 = p3.add_run()
r3.text = "We welcome questions from the jury and invite you to inspect our live digital twin & GCS demonstration!"
r3.font.name = "Times New Roman"
r3.font.size = Pt(11)
r3.font.italic = True
r3.font.color.rgb = COLOR_MUTED

print("Slide 14 configured.")

# Save presentation
prs.save(OUTPUT_PATH)
prs.save(DOWNLOADS_OUTPUT_PATH)
print(f"\n Presentation saved successfully to:")
print(f"1. {OUTPUT_PATH}")
print(f"2. {DOWNLOADS_OUTPUT_PATH}")
