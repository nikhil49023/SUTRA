import pptx

prs = pptx.Presentation('Smart_Horizon_2026_SUTRA_Grand_Finale_Pitch.pptx')

replacements = {
    "03 September 2026": "05 September 2026",
    "232 / 232 PyTests": "255 / 255 PyTests",
    "232 automated test suites": "255 automated test suites",
    "232-test verification suite": "255-test verification suite (100% Deterministic Pass)",
    "-5 dB SNR": "-8.0 dB SNR",
}

for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    for old_text, new_text in replacements.items():
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)

# Check Slide 14 and append GitHub link if not present
slide14 = prs.slides[13]
for shape in slide14.shapes:
    if shape.has_text_frame and "Grand Finals Project Team" in shape.text:
        if "https://github.com/nikhil49023/SUTRA" not in shape.text:
            p = shape.text_frame.add_paragraph()
            p.text = "Public GitHub Repository: https://github.com/nikhil49023/SUTRA"
            p.font.bold = True

prs.save('Smart_Horizon_2026_SUTRA_Grand_Finale_Pitch.pptx')
print("✅ Successfully updated Smart_Horizon_2026_SUTRA_Grand_Finale_Pitch.pptx!")
